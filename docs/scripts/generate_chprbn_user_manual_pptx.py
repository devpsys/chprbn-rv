#!/usr/bin/env python3
"""
Generate docs/CHPRBN_Field_App_User_Manual.pptx from a fixed slide outline.

Requirements:
    pip install python-pptx

Run from repository root:
    python docs/scripts/generate_chprbn_user_manual_pptx.py
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    raise SystemExit(
        "Missing dependency: install with  pip install python-pptx  then re-run."
    )

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "CHPRBN_Field_App_User_Manual.pptx"

# (title, bullet_lines, reserve_screenshot_space)
SLIDES = [
    (
        "CHPRBN Field Mobile App",
        [
            "User manual — field officers",
            "License lookup · Verification · Sync",
            "[ Add screenshot: app icon / splash ]",
        ],
    ),
    (
        "Overview",
        [
            "Adhoc login (username + password)",
            "QR or manual license lookup",
            "Local verification then upload (sync)",
            "[ Screenshot: Dashboard ]",
        ],
    ),
    (
        "Splash & session",
        [
            "Initializing ~2.5s then check cached user",
            "Valid token → Dashboard",
            "Invalid / missing → Login",
            "[ Screenshot: Splash ]",
        ],
    ),
    (
        "Login — success",
        [
            "adhoc/login + adhoc/profile",
            "User and token saved locally",
            "[ Screenshot: Login → Dashboard ]",
        ],
    ),
    (
        "Login — errors",
        [
            "“Username and password are required.”",
            "API message or “Login failed.” / profile errors",
            "Offline: “No cached session…” without valid cache",
            "[ Screenshot: inline error on Login ]",
        ],
    ),
    (
        "Dashboard",
        [
            "Loading · Success (tiles) · Error + Retry",
            "[ Screenshot: each state if possible ]",
        ],
    ),
    (
        "Scan QR",
        [
            "Camera permission · preview · torch",
            "Extract id from “#: …” line when present",
            "[ Screenshot: scan UI ]",
        ],
    ),
    (
        "Manual entry",
        [
            "Enter license number → Record detail",
            "[ Screenshot: manual entry ]",
        ],
    ),
    (
        "Record detail — states",
        [
            "Loading · Success · Not found · Error UI",
            "Not found: try again / enter manually",
            "Error UI: “Connection lost” + retry",
            "[ Screenshot: four states ]",
        ],
    ),
    (
        "Verification form",
        [
            "Location + remarks required",
            "“Mark Verified” must be ON to save",
            "Validation errors shown in red",
            "[ Screenshot: form + error example ]",
        ],
    ),
    (
        "Verified list",
        [
            "Filters: All, Active, Expired, Pending sync",
            "Row badges: Pending / Failed / Synced",
            "[ Screenshot: list with badges ]",
        ],
    ),
    (
        "Sync",
        [
            "Stats: Total, Synced, Pending, Failed",
            "Sync all · Retry failed",
            "Summary: “X ok, Y failed (Z attempted)”",
            "Load/sync errors in red text",
            "[ Screenshot: sync screen + summary ]",
        ],
    ),
    (
        "Profile & logout",
        [
            "Cached profile · Logout clears session → Login",
            "[ Screenshot: Profile ]",
        ],
    ),
    (
        "Screenshot placeholders",
        [
            "Duplicate this slide per feature if you need more space.",
            "[ Large empty area below for pasted screenshots ]",
        ],
    ),
]


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    try:
        blank = prs.slide_layouts[6]  # Blank
    except IndexError:
        blank = prs.slide_layouts[5]

    for title, bullets in SLIDES:
        slide = prs.slides.add_slide(blank)
        # Title
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        # Body
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(7.2), Inches(5.5))
        bt = body.text_frame
        bt.word_wrap = True
        for i, line in enumerate(bullets):
            para = bt.add_paragraph() if i > 0 else bt.paragraphs[0]
            para.text = line
            para.font.size = Pt(18)
            para.level = 0
        # Right placeholder for screenshot
        box = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(8.1),
            Inches(1.35),
            Inches(4.6),
            Inches(5.2),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(190, 190, 190)
        tf2 = box.text_frame
        tf2.text = "[ Screenshot ]"
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].font.italic = True

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
